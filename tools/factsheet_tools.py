"""
Herramientas para generar y actualizar Fact Sheets (.pptx) de los fondos.
Manipula archivos PowerPoint via python-pptx.

Fondos soportados (fondo_key):
  "PT"       → Toesca Rentas Inmobiliarias PT
  "Apoquindo"→ Toesca Rentas Inmobiliarias Apoquindo
  "TRI"      → Toesca Rentas Inmobiliarias (fondo madre)
"""
import json
import os
import re
import shutil
from datetime import date, datetime
from calendar import monthrange

from config import SHAREPOINT_DIR, WORK_DIR

# ── Paths dentro de SHAREPOINT_DIR ─────────────────────────────────────────
_INMOBILIARIO = "Inmobiliario Toesca - Documentos"

_FS_DIRS = {
    "PT":        os.path.join(_INMOBILIARIO, "Fondo Rentas PT", "Facts Sheet"),
    "Apoquindo": os.path.join(_INMOBILIARIO, "Fondo Rentas Apoquindo", "Facts Sheet"),
    "TRI":       os.path.join(_INMOBILIARIO, "Fondo Rentas Inmobiliarias TRI", "Facts Sheet"),
}

_FS_SUFFIXES = ("vActualizar.pptx", "vRevisar.pptx", "vF.pptx")

_FS_OUTPUT_BASE = {
    "PT":        "Fact Sheet - Toesca Rentas Inmobiliarias PT",
    "Apoquindo": "Fact Sheet - Toesca Rentas Inmobiliarias Apoquindo",
    "TRI":       "Fact Sheet - Toesca Rentas Inmobiliarias",
}

_MESES_ES = {
    1: "Enero", 2: "Febrero", 3: "Marzo", 4: "Abril",
    5: "Mayo", 6: "Junio", 7: "Julio", 8: "Agosto",
    9: "Septiembre", 10: "Octubre", 11: "Noviembre", 12: "Diciembre",
}


# ── Helpers internos ────────────────────────────────────────────────────────

def _sp(subfolder: str) -> str:
    """Devuelve ruta absoluta dentro de SHAREPOINT_DIR."""
    return os.path.join(SHAREPOINT_DIR, subfolder)


def _fs_folder(fondo_key: str, año: int | None = None, mes: int | None = None) -> str:
    """Ruta a la carpeta de Fact Sheets del fondo."""
    base = _sp(_FS_DIRS[fondo_key])
    if not año:
        return base
    if fondo_key == "PT":
        # PT usa solo subcarpeta YYYY (sin mes)
        return os.path.join(base, str(año))
    # Apoquindo y TRI usan subcarpetas YYYY/MesNombre
    if mes:
        return os.path.join(base, str(año), _MESES_ES[mes])
    return os.path.join(base, str(año))


def _find_template(fondo_key: str, año: int, mes: int) -> str | None:
    """Busca el archivo de trabajo (vActualizar, vRevisar o vF) en la carpeta del fondo."""
    folder = _fs_folder(fondo_key, año, mes)
    if not os.path.isdir(folder):
        return None
    for name in os.listdir(folder):
        if any(name.endswith(s) for s in _FS_SUFFIXES):
            return os.path.join(folder, name)
    return None


def _work_path(fondo_key: str) -> str:
    """Ruta del archivo de trabajo en WORK_DIR."""
    return os.path.join(WORK_DIR, f"_fs_{fondo_key}_trabajo.pptx")


def _get_shape(slide, name: str):
    """Devuelve el shape con ese nombre en el slide, o None."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _set_cell_text(cell, text: str):
    """Actualiza el texto de una celda preservando el formato del primer run."""
    tf = cell.text_frame
    if not tf.paragraphs:
        return
    para = tf.paragraphs[0]
    if para.runs:
        # Actualiza primer run; borra los demás si hay más de uno
        para.runs[0].text = str(text)
        for run in para.runs[1:]:
            run.text = ""
    else:
        # Sin runs: crear uno via XML
        from lxml import etree
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        r = etree.SubElement(para._p, f"{{{ns}}}r")
        t = etree.SubElement(r, f"{{{ns}}}t")
        t.text = str(text)


def _update_table(slide, shape_name: str, updates: dict) -> bool:
    """
    Actualiza celdas de una tabla.
    updates = {(row_idx, col_idx): "nuevo texto", ...}
    Devuelve True si encontró la tabla.
    """
    shape = _get_shape(slide, shape_name)
    if shape is None or not shape.has_table:
        return False
    tbl = shape.table
    for (r, c), text in updates.items():
        if r < len(tbl.rows) and c < len(tbl.rows[r].cells):
            _set_cell_text(tbl.rows[r].cells[c], text)
    return True


# ── Herramientas públicas ───────────────────────────────────────────────────

# ─── Reglas de fecha FS ──────────────────────────────────────────────────────
# Los FS se publican en enero, abril, julio y octubre.
# El parámetro `mes` en todas las funciones de FS es el mes de publicación.
#   Bursátil: 3 meses terminando en el mes del FS  (ej. enero → nov/dic/ene)
#   Contable: último día del trimestre cerrado anterior al mes del FS
#             enero→dic, abril→mar, julio→jun, octubre→sep
# ─────────────────────────────────────────────────────────────────────────────

_FS_MES_A_CONTABLE = {1: 12, 4: 3, 7: 6, 10: 9}


def fecha_contable_fs(año: int, mes_fs: int) -> date:
    """
    Retorna la fecha de cierre contable para el mes del FS dado.
      enero(1)   → 31-dic-año anterior
      abril(4)   → 31-mar-año actual
      julio(7)   → 30-jun-año actual
      octubre(10)→ 30-sep-año actual
    """
    mes_cont = _FS_MES_A_CONTABLE[mes_fs]
    año_cont = año - 1 if mes_fs == 1 else año
    ultimo_dia = monthrange(año_cont, mes_cont)[1]
    return date(año_cont, mes_cont, ultimo_dia)


def _fmt_clp(valor: float) -> str:
    """12000.0 → '$12.000'  |  1234567.0 → '$1.234.567'"""
    return "$" + f"{round(valor):,}".replace(",", ".")


def _mes_anterior(año: int, mes: int, n: int = 1):
    """Retorna (año, mes) n meses antes."""
    m = mes - n
    a = año
    while m < 1:
        m += 12
        a -= 1
    return a, m


def _parse_precio_cuota(texto: str):
    """
    Parsea el string retornado por obtener_precio_cuota.
    Retorna {"fecha": "31-10-2025", "valor": "$12.000"} o None si hay error.
    """
    m = re.search(r"al (\d{2}/\d{2}/\d{4}).*?:\s*([\d,.]+)", texto)
    if not m:
        return None
    fecha = m.group(1).replace("/", "-")      # "31/10/2025" → "31-10-2025"
    precio = float(m.group(2).replace(",", ""))
    return {"fecha": fecha, "valor": _fmt_clp(precio)}


def obtener_valor_libro_fs(fondo_key: str, año_fs: int, mes_fs: int) -> str:
    """
    Retorna el valor cuota libro formateado para la tabla 'EL FONDO' del Fact Sheet.
    Usa la fecha contable del FS (última del trimestre cerrado anterior al mes del FS).

    fondo_key: 'A&R PT', 'A&R Apoquindo', 'A&R Rentas'
    mes_fs: 1, 4, 7 ó 10

    Retorna texto con fecha, valor y JSON listo para datos_json['info_fondo'].
    Para A&R Rentas retorna las 3 series (A, C, I).
    """
    from tools.eeff_tools import extraer_datos_eeff, buscar_pdf_eeff

    fc = fecha_contable_fs(año_fs, mes_fs)
    año_cont, mes_cont = fc.year, fc.month
    fecha_str = fc.strftime("%d-%m-%Y")

    pdf_path = buscar_pdf_eeff(fondo_key, año_cont, mes_cont)
    if not os.path.isfile(pdf_path):
        return f"No se encontró el PDF de EEFF para {fondo_key} {mes_cont}/{año_cont}.\n{pdf_path}"

    datos = extraer_datos_eeff(pdf_path, fondo_key)

    if datos.get("error"):
        return f"Error al leer EEFF {fondo_key}: {datos['error']}"

    if not datos["valor_cuota"]:
        return f"No se pudo extraer valor cuota de {fondo_key} {mes_cont}/{año_cont}. Revisar PDF manualmente."

    lineas = [f"Valor cuota libro {fondo_key} al {fecha_str}:"]
    info_fondo_json = {"valor_libro_fecha": fecha_str}

    if fondo_key == "A&R Rentas":
        # Multi-serie A, C, I
        for serie, val in sorted(datos["valor_cuota"].items(), key=lambda x: str(x[0])):
            label = f"Serie {serie}" if serie else "Único"
            lineas.append(f"  {label}: {_fmt_clp(val)}")
        # Para FS TRI pasamos las 3 series
        info_fondo_json["series"] = {
            s: _fmt_clp(v) for s, v in datos["valor_cuota"].items()
        }
    else:
        # Fondo único (PT o Apoquindo)
        val = list(datos["valor_cuota"].values())[0]
        lineas.append(f"  {_fmt_clp(val)}")
        info_fondo_json["valor_libro"] = _fmt_clp(val)

    lineas.append(f"\nJSON listo para datos_json['info_fondo']:\n{json.dumps(info_fondo_json, ensure_ascii=False)}")
    return "\n".join(lineas)


def obtener_precios_bursatiles_fs(nemotecnico: str, año: int, mes: int, n: int = 3) -> str:
    """
    Retorna los últimos n meses de precios de cuota formateados para el Fact Sheet.
    Llama a obtener_precio_cuota para cada mes y parsea el resultado.

    Retorna JSON con lista de {"fecha": "DD-MM-YYYY", "valor": "$XX.XXX"},
    ordenada del más antiguo al más reciente.

    Ejemplo de uso:
        obtener_precios_bursatiles_fs('CFITRIPT-E', 2026, 1, n=3)
        → '[{"fecha": "30-11-2025", "valor": "$12.000"}, ...]'
    """
    from tools.web_bursatil_tools import obtener_precio_cuota

    meses = [_mes_anterior(año, mes, i) for i in range(n - 1, -1, -1)]
    resultados = []
    errores = []

    for a, m in meses:
        texto = obtener_precio_cuota(nemotecnico, a, m)
        parsed = _parse_precio_cuota(texto)
        if parsed:
            resultados.append(parsed)
        else:
            errores.append(f"{m:02d}/{a}: {texto}")

    salida = f"Precios bursátiles {nemotecnico} (últimos {n} meses):\n"
    for r in resultados:
        salida += f"  {r['fecha']}  {r['valor']}\n"
    if errores:
        salida += "\nErrores:\n" + "\n".join(f"  {e}" for e in errores)
    salida += f"\nJSON listo para datos_json['precios_bursatiles']:\n{json.dumps(resultados, ensure_ascii=False)}"
    return salida


def listar_shapes_fs(fondo_key: str, año: int, mes: int) -> str:
    """
    Lista todos los shapes del Slide 1 del Fact Sheet del fondo indicado.
    Útil para descubrir nombres de tablas y text boxes.
    fondo_key: "PT", "Apoquindo" o "TRI"
    """
    if fondo_key not in _FS_DIRS:
        return f"fondo_key inválido: {fondo_key}. Usar 'PT', 'Apoquindo' o 'TRI'."

    try:
        from pptx import Presentation
    except ImportError:
        return "Error: python-pptx no está instalado. Ejecutar: pip install python-pptx"

    src = _find_template(fondo_key, año, mes)
    if not src:
        folder = _fs_folder(fondo_key, año, mes)
        return f"No se encontró archivo vActualizar/vRevisar en:\n{folder}"

    prs = Presentation(src)
    slide = prs.slides[0]
    lines = [f"Fact Sheet {fondo_key} — Shapes en Slide 1 ({src}):\n"]
    for shape in slide.shapes:
        tipo = shape.shape_type
        info = f"  [{shape.name}]  type={tipo}"
        if shape.has_table:
            tbl = shape.table
            info += f"  TABLE {len(tbl.rows)}x{len(tbl.columns)}"
            # Mostrar primera fila (encabezado)
            header = [cell.text_frame.text.strip()[:30] for cell in tbl.rows[0].cells]
            info += f"\n      header: {header}"
        elif shape.has_text_frame:
            preview = shape.text_frame.text.strip()[:60].replace("\n", " ")
            info += f"  TEXT: {repr(preview)}"
        lines.append(info)
    return "\n".join(lines)


def leer_tabla_fs(fondo_key: str, año: int, mes: int, shape_name: str) -> str:
    """
    Lee el contenido de una tabla específica del Fact Sheet (Slide 1).
    Útil para inspeccionar datos actuales antes de actualizar.
    """
    if fondo_key not in _FS_DIRS:
        return f"fondo_key inválido: {fondo_key}."

    try:
        from pptx import Presentation
    except ImportError:
        return "Error: python-pptx no instalado."

    # Intentar primero desde WORK_DIR (archivo de trabajo), si no desde SharePoint
    work = _work_path(fondo_key)
    src = work if os.path.exists(work) else _find_template(fondo_key, año, mes)
    if not src:
        return f"No se encontró archivo FS para {fondo_key} {mes}/{año}."

    prs = Presentation(src)
    slide = prs.slides[0]
    shape = _get_shape(slide, shape_name)
    if shape is None:
        return f"No se encontró shape '{shape_name}' en Slide 1."
    if not shape.has_table:
        return f"'{shape_name}' no es una tabla."

    tbl = shape.table
    lines = [f"Tabla '{shape_name}':"]
    for r, row in enumerate(tbl.rows):
        cells = [cell.text_frame.text.strip() for cell in row.cells]
        lines.append(f"  Fila {r}: {cells}")
    return "\n".join(lines)


def preparar_fs(fondo_key: str, año: int, mes: int) -> str:
    """
    Copia el archivo vActualizar/vRevisar del fondo desde SharePoint a WORK_DIR
    para comenzar la edición. Devuelve confirmación con la ruta de trabajo.
    Siempre llamar esto antes de actualizar_fs_*.
    """
    if fondo_key not in _FS_DIRS:
        return f"fondo_key inválido: {fondo_key}."

    src = _find_template(fondo_key, año, mes)
    if not src:
        folder = _fs_folder(fondo_key, año, mes)
        return (
            f"No se encontró archivo de plantilla en:\n{folder}\n"
            f"Verifica que el archivo vActualizar/vRevisar exista para {_MESES_ES.get(mes, mes)} {año}."
        )

    os.makedirs(WORK_DIR, exist_ok=True)
    dest = _work_path(fondo_key)
    shutil.copy2(src, dest)
    return (
        f"FS {fondo_key} copiado a WORK_DIR.\n"
        f"  Origen:  {src}\n"
        f"  Trabajo: {dest}\n"
        f"Listo para actualizar tablas."
    )


def actualizar_fs_pt(año: int, mes: int, datos_json: str) -> str:
    """
    Actualiza todas las tablas numéricas del Slide 1 del Fact Sheet PT.
    Requiere haber llamado preparar_fs('PT', año, mes) antes.

    datos_json es un JSON con la siguiente estructura (todos los campos son opcionales;
    solo se actualizan los que se incluyan):

    {
      "precios_bursatiles": [
        {"fecha": "31-10-2025", "valor": "$12.000"},
        {"fecha": "31-11-2025", "valor": "$12.000"},
        {"fecha": "31-12-2025", "valor": "$12.000"}
      ],
      "valor_libro": [
        {"fecha": "30-06-2025", "valor": "$12.261"},
        {"fecha": "30-09-2025", "valor": "$13.184"},
        {"fecha": "31-12-2025", "valor": "$13.707"}
      ],
      "rentabilidad": {
        "inicio_bursatil": "-6,4%", "ytd_bursatil": "0,1%", "12m_bursatil": "8,6%",
        "inicio_libro":    "-6,2%", "ytd_libro":    "10,7%","12m_libro":    "20,8%",
        "dy_bursatil": "9,1%", "dy_libro": "8,2%",
        "dy_amort_bursatil": "9,1%", "dy_amort_libro": "8,2%"
      },
      "dividendos": [
        {"fecha": "29-04-2025", "concepto": "Dividendo Provisorio", "monto": "$440,8"},
        {"fecha": "28-07-2025", "concepto": "Dividendo Provisorio", "monto": "$299,0"},
        {"fecha": "22-10-2025", "concepto": "Dividendo Provisorio", "monto": "$337,5"},
        {"fecha": "29-12-2025", "concepto": "Dividendo Provisorio", "monto": "$339,1"}
      ],
      "otros_indicadores": {
        "tasa_arriendo_bursatil": "7,3%", "cap_rate_bursatil": "6,0%",
        "tasa_arriendo_libro":    "7,4%", "cap_rate_libro":    "6,0%",
        "ingresos_u12m": "UF 210.966", "ingresos_mes": "UF 17.483",
        "noi_u12m": "UF 171.668",      "noi_mes": "UF 15.067"
      },
      "balance": {
        "fecha": "31 de diciembre de 2025",
        "efectivo": "487.675",
        "otros_activos_corrientes": "355.335",
        "propiedades_inversion": "115.816.859",
        "activo_impuestos_diferidos": "10.867.793",
        "total_activos": "127.527.662",
        "prestamos_bancarios": "90.747.171",
        "pasivos_impuestos_diferidos": "10.484.360",
        "otros_pasivos": "3.815.939",
        "patrimonio": "22.480.192",
        "total_pasivos_patrimonio": "127.527.662"
      },
      "gastos": {
        "fecha": "30 DE SEPTIEMBRE DE 2025",
        "comision": "69.857",
        "recurrentes": "38.169",
        "otros": "343",
        "total": "108.360"
      },
      "endeudamiento": {
        "leverage": "4,35 x",
        "ltv": "81,2%",
        "tasa_promedio": "4,2%",
        "duration": "4,4",
        "deuda_neta": "2.346.979"
      },
      "perfil_vencimiento": {
        "0_3": "3%", "3_7": "97%", "7_10": "0%", "mas_10": "0%"
      },
      "info_fondo": {
        "valor_libro_fecha": "30-09-2025",
        "valor_libro": "$13.184"
      }
    }
    """
    try:
        from pptx import Presentation
    except ImportError:
        return "Error: python-pptx no instalado."

    work = _work_path("PT")
    if not os.path.exists(work):
        return "No hay archivo de trabajo. Llama primero a preparar_fs('PT', año, mes)."

    try:
        datos = json.loads(datos_json)
    except json.JSONDecodeError as e:
        return f"Error al parsear datos_json: {e}"

    prs = Presentation(work)
    slide = prs.slides[0]
    actualizados = []
    omitidos = []

    # ── Fecha del mes (text box) ──────────────────────────────────────────
    mes_str = f"{_MESES_ES[mes].upper()} {año}"
    updated_fecha = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            for m_es in _MESES_ES.values():
                if m_es.upper() in txt and any(str(y) in txt for y in range(2020, 2035)):
                    para = shape.text_frame.paragraphs[0]
                    if para.runs:
                        para.runs[0].text = mes_str
                        for run in para.runs[1:]:
                            run.text = ""
                    else:
                        from lxml import etree
                        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
                        r = etree.SubElement(para._p, f"{{{ns}}}r")
                        t = etree.SubElement(r, f"{{{ns}}}t")
                        t.text = mes_str
                    updated_fecha = True
                    break
        if updated_fecha:
            break
    if updated_fecha:
        actualizados.append(f"Fecha → {mes_str}")
    else:
        omitidos.append("Fecha del mes (text box no encontrado)")

    # ── Tabla 2: Valor cuota bursátil (últimos 3 meses) ──────────────────
    pb = datos.get("precios_bursatiles", [])
    if pb:
        updates = {}
        for i, item in enumerate(pb[:3]):
            row = 2 + i  # filas 2, 3, 4
            updates[(row, 0)] = item.get("fecha", "")
            updates[(row, 1)] = item.get("valor", "")
        if _update_table(slide, "Tabla 2", updates):
            actualizados.append("Tabla 2 (precios bursátiles)")
        else:
            omitidos.append("Tabla 2 (shape no encontrado)")

    # ── Tabla 7: Valor cuota libro (últimos 3 trimestres) ─────────────────
    vl = datos.get("valor_libro", [])
    if vl:
        updates = {}
        for i, item in enumerate(vl[:3]):
            row = 2 + i
            updates[(row, 0)] = item.get("fecha", "")
            updates[(row, 1)] = item.get("valor", "")
        if _update_table(slide, "Tabla 7", updates):
            actualizados.append("Tabla 7 (valor cuota libro)")
        else:
            omitidos.append("Tabla 7 (shape no encontrado)")

    # ── Tabla 52: Rentabilidad ─────────────────────────────────────────────
    rent = datos.get("rentabilidad", {})
    if rent:
        updates = {
            (2, 1): rent.get("inicio_bursatil", ""),
            (2, 2): rent.get("inicio_libro", ""),
            (3, 1): rent.get("ytd_bursatil", ""),
            (3, 2): rent.get("ytd_libro", ""),
            (4, 1): rent.get("12m_bursatil", ""),
            (4, 2): rent.get("12m_libro", ""),
            (5, 1): rent.get("dy_bursatil", ""),
            (5, 2): rent.get("dy_libro", ""),
            (6, 1): rent.get("dy_amort_bursatil", ""),
            (6, 2): rent.get("dy_amort_libro", ""),
        }
        updates = {k: v for k, v in updates.items() if v}
        if _update_table(slide, "Tabla 52", updates):
            actualizados.append("Tabla 52 (rentabilidad)")
        else:
            omitidos.append("Tabla 52 (shape no encontrado)")

    # ── Tabla 50: Dividendos (repartos últimos 12 meses) ──────────────────
    divs = datos.get("dividendos", [])
    if divs:
        updates = {}
        for i, div in enumerate(divs[:4]):
            row = 2 + i
            updates[(row, 0)] = div.get("fecha", "")
            updates[(row, 1)] = div.get("concepto", "")
            updates[(row, 2)] = div.get("monto", "")
        if _update_table(slide, "Tabla 50", updates):
            actualizados.append("Tabla 50 (dividendos)")
        else:
            omitidos.append("Tabla 50 (shape no encontrado)")

    # ── Tabla 6: Otros indicadores ────────────────────────────────────────
    noi = datos.get("otros_indicadores", {})
    if noi:
        updates = {
            (1, 1): noi.get("tasa_arriendo_bursatil", ""),
            (1, 4): noi.get("tasa_arriendo_libro", ""),
            (2, 1): noi.get("cap_rate_bursatil", ""),
            (2, 4): noi.get("cap_rate_libro", ""),
            (3, 1): noi.get("ingresos_u12m", ""),
            (4, 1): noi.get("ingresos_mes", ""),
            (5, 1): noi.get("noi_u12m", ""),
            (6, 1): noi.get("noi_mes", ""),
        }
        updates = {k: v for k, v in updates.items() if v}
        if _update_table(slide, "Tabla 6", updates):
            actualizados.append("Tabla 6 (otros indicadores)")
        else:
            omitidos.append("Tabla 6 (shape no encontrado)")

    # ── Tabla 4: Balance consolidado ──────────────────────────────────────
    bal = datos.get("balance", {})
    if bal:
        # Fila 0 = encabezado con fecha
        fecha_bal = bal.get("fecha", "")
        header_txt = f"BALANCE CONSOLIDADO AL {fecha_bal.upper()} (en miles de pesos)" if fecha_bal else ""
        updates = {}
        if header_txt:
            updates[(0, 0)] = header_txt
        updates.update({
            (1, 0): "Efectivo y Efectivo Equivalente",
            (1, 1): bal.get("efectivo", ""),
            (1, 3): "Préstamos Bancarios",
            (1, 4): bal.get("prestamos_bancarios", ""),
            (2, 1): bal.get("otros_activos_corrientes", ""),
            (2, 4): bal.get("pasivos_impuestos_diferidos", ""),
            (3, 1): bal.get("propiedades_inversion", ""),
            (3, 4): bal.get("otros_pasivos", ""),
            (4, 1): bal.get("activo_impuestos_diferidos", ""),
            (4, 4): bal.get("patrimonio", ""),
            (5, 1): bal.get("total_activos", ""),
            (5, 4): bal.get("total_pasivos_patrimonio", ""),
        })
        updates = {k: v for k, v in updates.items() if v}
        if _update_table(slide, "Tabla 4", updates):
            actualizados.append("Tabla 4 (balance consolidado)")
        else:
            omitidos.append("Tabla 4 (shape no encontrado)")

    # ── Tabla 44: Gastos del fondo ────────────────────────────────────────
    gas = datos.get("gastos", {})
    if gas:
        fecha_gas = gas.get("fecha", "")
        header_txt = f"GASTOS DEL FONDO {fecha_gas.upper()} (en miles de pesos)" if fecha_gas else ""
        updates = {}
        if header_txt:
            updates[(0, 0)] = header_txt
        updates.update({
            (1, 1): gas.get("comision", ""),
            (2, 1): gas.get("recurrentes", ""),
            (3, 1): gas.get("otros", ""),
            (4, 1): gas.get("total", ""),
        })
        updates = {k: v for k, v in updates.items() if v}
        if _update_table(slide, "Tabla 44", updates):
            actualizados.append("Tabla 44 (gastos)")
        else:
            omitidos.append("Tabla 44 (shape no encontrado)")

    # ── Tabla 8: Endeudamiento ────────────────────────────────────────────
    end = datos.get("endeudamiento", {})
    if end:
        updates = {
            (1, 1): end.get("leverage", ""),
            (2, 1): end.get("ltv", ""),
            (3, 1): end.get("tasa_promedio", ""),
            (4, 1): end.get("duration", ""),
            (5, 1): end.get("deuda_neta", ""),
        }
        updates = {k: v for k, v in updates.items() if v}
        if _update_table(slide, "Tabla 8", updates):
            actualizados.append("Tabla 8 (endeudamiento)")
        else:
            omitidos.append("Tabla 8 (shape no encontrado)")

    # ── Tabla 3: Perfil de vencimiento deuda ──────────────────────────────
    perf = datos.get("perfil_vencimiento", {})
    if perf:
        updates = {
            (1, 1): perf.get("0_3", ""),
            (2, 1): perf.get("3_7", ""),
            (3, 1): perf.get("7_10", ""),
            (4, 1): perf.get("mas_10", ""),
        }
        updates = {k: v for k, v in updates.items() if v}
        if _update_table(slide, "Tabla 3", updates):
            actualizados.append("Tabla 3 (perfil vencimiento deuda)")
        else:
            omitidos.append("Tabla 3 (shape no encontrado)")

    # ── Tabla 19: Info del fondo (valor libro reciente) ───────────────────
    info = datos.get("info_fondo", {})
    if info:
        updates = {}
        if info.get("valor_libro_fecha"):
            updates[(3, 0)] = f"Valor Libro  {info['valor_libro_fecha']}"
        if info.get("valor_libro"):
            updates[(3, 1)] = info["valor_libro"]
        if updates and _update_table(slide, "Tabla 19", updates):
            actualizados.append("Tabla 19 (info fondo)")
        elif updates:
            omitidos.append("Tabla 19 (shape no encontrado)")

    prs.save(work)

    resultado = f"FS PT actualizado en WORK_DIR.\n"
    if actualizados:
        resultado += f"\nActualizados ({len(actualizados)}):\n" + "\n".join(f"  ✓ {x}" for x in actualizados)
    if omitidos:
        resultado += f"\nOmitidos (no encontrados):\n" + "\n".join(f"  ✗ {x}" for x in omitidos)
    resultado += f"\n\nArchivo de trabajo: {work}"
    resultado += "\n\nLlama a guardar_fs('PT', año, mes) para guardarlo en SharePoint."
    return resultado


def actualizar_fs_apoquindo(año: int, mes: int, datos_json: str) -> str:
    """
    Actualiza las tablas numéricas del Slide 1 del Fact Sheet Apoquindo.
    Requiere haber llamado preparar_fs('Apoquindo', año, mes) antes.
    Apoquindo no tiene ticker bursátil ni tabla de repartos.

    datos_json estructura:
    {
      "valor_libro": [
        {"fecha": "30-09-2025", "valor": "$14.500"},
        {"fecha": "31-12-2025", "valor": "$15.200"},
        {"fecha": "31-03-2026", "valor": "$15.800"}
      ],
      "rentabilidad": {
        "inicio_libro": "-3,0%", "ytd_libro": "7,4%", "12m_libro": "7,4%",
        "dy_libro": "0,0%", "dy_amort_libro": "0,0%"
      },
      "otros_indicadores": {
        "tasa_arriendo": "5,6%", "cap_rate": "4,8%",
        "noi_u12m": "UF 201.789", "noi_mes": "UF 17.556",
        "ingresos_u12m": "UF 172.002", "ingresos_mes": "UF 15.107"
      },
      "gastos": {
        "fecha": "31 DICIEMBRE 2025",
        "comision": "145.795", "recurrentes": "38.774", "otros": "", "total": "184.569"
      },
      "balance": {
        "fecha": "31 de diciembre de 2025",
        "efectivo": "...", "otros_activos_corrientes": "...",
        "propiedades_inversion": "...", "activo_impuestos_diferidos": "...", "total_activos": "...",
        "prestamos_bancarios": "...", "pasivos_impuestos_diferidos": "...",
        "otros_pasivos": "...", "patrimonio": "...", "total_pasivos_patrimonio": "..."
      },
      "endeudamiento": {
        "leverage": "...", "ltv": "...", "tasa_promedio": "...", "duration": "...", "deuda_neta": "..."
      },
      "perfil_vencimiento": {"0_3": "...", "3_7": "...", "7_10": "...", "mas_10": "..."},
      "info_fondo": {"valor_libro_fecha": "31-12-2025", "valor_libro": "$15.200"}
    }
    """
    try:
        from pptx import Presentation
    except ImportError:
        return "Error: python-pptx no instalado."

    work = _work_path("Apoquindo")
    if not os.path.exists(work):
        return "No hay archivo de trabajo. Llama primero a preparar_fs('Apoquindo', año, mes)."

    try:
        datos = json.loads(datos_json)
    except json.JSONDecodeError as e:
        return f"Error al parsear datos_json: {e}"

    prs = Presentation(work)
    slide = prs.slides[0]
    actualizados = []
    omitidos = []

    # Fecha del mes
    mes_str = f"{_MESES_ES[mes].upper()} {año}"
    updated_fecha = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            for m_es in _MESES_ES.values():
                if m_es.upper() in txt and any(str(y) in txt for y in range(2020, 2035)):
                    para = shape.text_frame.paragraphs[0]
                    if para.runs:
                        para.runs[0].text = mes_str
                        for run in para.runs[1:]:
                            run.text = ""
                    updated_fecha = True
                    break
        if updated_fecha:
            break
    if updated_fecha:
        actualizados.append(f"Fecha → {mes_str}")

    # Tabla 7: Valor cuota libro (3 trimestres)
    vl = datos.get("valor_libro", [])
    if vl:
        updates = {}
        for i, item in enumerate(vl[:3]):
            row = 2 + i
            updates[(row, 0)] = item.get("fecha", "")
            updates[(row, 1)] = item.get("valor", "")
        if _update_table(slide, "Tabla 7", {k: v for k, v in updates.items() if v}):
            actualizados.append("Tabla 7 (valor cuota libro)")
        else:
            omitidos.append("Tabla 7")

    # Tabla 46: Rentabilidad (solo columna Libro; Bursátil = N/A)
    rent = datos.get("rentabilidad", {})
    if rent:
        updates = {
            (2, 2): rent.get("inicio_libro", ""),
            (3, 2): rent.get("ytd_libro", ""),
            (4, 2): rent.get("12m_libro", ""),
            (5, 2): rent.get("dy_libro", ""),
            (6, 2): rent.get("dy_amort_libro", ""),
        }
        if _update_table(slide, "Tabla 46", {k: v for k, v in updates.items() if v}):
            actualizados.append("Tabla 46 (rentabilidad)")
        else:
            omitidos.append("Tabla 46")

    # Tabla 45: Otros indicadores (solo 2 cols, sin split bursátil/libro)
    noi = datos.get("otros_indicadores", {})
    if noi:
        updates = {
            (1, 1): noi.get("tasa_arriendo", ""),
            (2, 1): noi.get("cap_rate", ""),
            (3, 1): noi.get("noi_u12m", ""),
            (4, 1): noi.get("noi_mes", ""),
            (5, 1): noi.get("ingresos_u12m", ""),
            (6, 1): noi.get("ingresos_mes", ""),
        }
        if _update_table(slide, "Tabla 45", {k: v for k, v in updates.items() if v}):
            actualizados.append("Tabla 45 (otros indicadores)")
        else:
            omitidos.append("Tabla 45")

    # Tabla 50: Gastos del fondo
    gas = datos.get("gastos", {})
    if gas:
        fecha_gas = gas.get("fecha", "")
        updates = {}
        if fecha_gas:
            updates[(0, 0)] = f"GASTOS DEL FONDO AL {fecha_gas.upper()} (en miles de pesos)"
        updates.update({
            (1, 1): gas.get("comision", ""),
            (2, 1): gas.get("recurrentes", ""),
            (3, 1): gas.get("otros", ""),
            (4, 1): gas.get("total", ""),
        })
        if _update_table(slide, "Tabla 50", {k: v for k, v in updates.items() if v}):
            actualizados.append("Tabla 50 (gastos)")
        else:
            omitidos.append("Tabla 50")

    # Tabla 4: Balance (misma estructura que PT)
    bal = datos.get("balance", {})
    if bal:
        fecha_bal = bal.get("fecha", "")
        updates = {}
        if fecha_bal:
            updates[(0, 0)] = f"BALANCE CONSOLIDADO AL {fecha_bal.upper()} (en miles de pesos)"
        updates.update({
            (1, 1): bal.get("efectivo", ""),
            (1, 4): bal.get("prestamos_bancarios", ""),
            (2, 1): bal.get("otros_activos_corrientes", ""),
            (2, 4): bal.get("pasivos_impuestos_diferidos", ""),
            (3, 1): bal.get("propiedades_inversion", ""),
            (3, 4): bal.get("otros_pasivos", ""),
            (4, 1): bal.get("activo_impuestos_diferidos", ""),
            (4, 4): bal.get("patrimonio", ""),
            (5, 1): bal.get("total_activos", ""),
            (5, 4): bal.get("total_pasivos_patrimonio", ""),
        })
        if _update_table(slide, "Tabla 4", {k: v for k, v in updates.items() if v}):
            actualizados.append("Tabla 4 (balance)")
        else:
            omitidos.append("Tabla 4")

    # Tabla 48: Endeudamiento
    end = datos.get("endeudamiento", {})
    if end:
        updates = {
            (1, 1): end.get("leverage", ""),
            (2, 1): end.get("ltv", ""),
            (3, 1): end.get("tasa_promedio", ""),
            (4, 1): end.get("duration", ""),
            (5, 1): end.get("deuda_neta", ""),
        }
        if _update_table(slide, "Tabla 48", {k: v for k, v in updates.items() if v}):
            actualizados.append("Tabla 48 (endeudamiento)")
        else:
            omitidos.append("Tabla 48")

    # Tabla 3: Perfil vencimiento
    perf = datos.get("perfil_vencimiento", {})
    if perf:
        updates = {
            (1, 1): perf.get("0_3", ""),
            (2, 1): perf.get("3_7", ""),
            (3, 1): perf.get("7_10", ""),
            (4, 1): perf.get("mas_10", ""),
        }
        if _update_table(slide, "Tabla 3", {k: v for k, v in updates.items() if v}):
            actualizados.append("Tabla 3 (perfil vencimiento)")
        else:
            omitidos.append("Tabla 3")

    # Tabla 19: Info fondo
    info = datos.get("info_fondo", {})
    if info:
        updates = {}
        if info.get("valor_libro_fecha"):
            updates[(3, 0)] = f"Valor Libro  {info['valor_libro_fecha']}"
        if info.get("valor_libro"):
            updates[(3, 1)] = info["valor_libro"]
        if updates and _update_table(slide, "Tabla 19", updates):
            actualizados.append("Tabla 19 (info fondo)")

    prs.save(work)

    resultado = f"FS Apoquindo actualizado en WORK_DIR.\n"
    if actualizados:
        resultado += f"\nActualizados ({len(actualizados)}):\n" + "\n".join(f"  ✓ {x}" for x in actualizados)
    if omitidos:
        resultado += f"\nOmitidos:\n" + "\n".join(f"  ✗ {x}" for x in omitidos)
    resultado += f"\n\nLlama a guardar_fs('Apoquindo', año, mes) para guardar en SharePoint."
    return resultado


def actualizar_fs_tri(año: int, mes: int, datos_json: str) -> str:
    """
    Actualiza las tablas numéricas del Slide 1 del Fact Sheet TRI (Rentas Inmobiliarias).
    Requiere haber llamado preparar_fs('TRI', año, mes) antes.
    TRI tiene 3 series (A, C, I) en rentabilidad, precios bursátiles, libro y repartos.

    datos_json estructura:
    {
      "precios_bursatiles": [
        {"fecha": "30-11-2025", "serie_a": "$16.490", "serie_c": "$16.750", "serie_i": "$29.476"},
        {"fecha": "31-12-2025", "serie_a": "$15.309", "serie_c": "$16.750", "serie_i": "$29.476"},
        {"fecha": "31-01-2026", "serie_a": "$15.309", "serie_c": "$16.750", "serie_i": "$29.476"}
      ],
      "valor_libro": [
        {"fecha": "31-03-2025", "serie_a": "$29.591", "serie_c": "$29.994", "serie_i": "$30.141"},
        {"fecha": "30-06-2025", "serie_a": "$29.989", "serie_c": "$30.393", "serie_i": "$30.539"},
        {"fecha": "30-09-2025", "serie_a": "$31.775", "serie_c": "$32.178", "serie_i": "$32.324"},
        {"fecha": "31-12-2025", "serie_a": "$31.869", "serie_c": "$32.252", "serie_i": "$32.390"}
      ],
      "rentabilidad": {
        "inicio_bursatil_a": "-8,4%", "inicio_libro_a": "1,7%",
        "inicio_bursatil_c": "-6,2%", "inicio_libro_c": "1,5%",
        "inicio_bursatil_i": "-0,9%", "inicio_libro_i": "2,0%",
        "ytd_bursatil_a": "-0,2%",    "ytd_libro_a": "22,1%",
        "ytd_bursatil_c": "3,1%",     "ytd_libro_c": "15,2%",
        "ytd_bursatil_i": "-0,7%",    "ytd_libro_i": "18,8%",
        "12m_bursatil_a": "-0,2%",    "12m_libro_a": "22,1%",
        "12m_bursatil_c": "3,1%",     "12m_libro_c": "15,2%",
        "12m_bursatil_i": "-0,7%",    "12m_libro_i": "18,8%",
        "dy_bursatil_a": "4,1%",      "dy_libro_a": "2,2%",
        "dy_bursatil_c": "4,6%",      "dy_libro_c": "2,4%",
        "dy_bursatil_i": "2,8%",      "dy_libro_i": "2,5%",
        "dy_amort_bursatil_a": "34,7%","dy_amort_libro_a": "18,3%",
        "dy_amort_bursatil_c": "35,4%","dy_amort_libro_c": "18,3%",
        "dy_amort_bursatil_i": "20,2%","dy_amort_libro_i": "16,4%"
      },
      "dividendos": [
        {"fecha": "29-04-2025", "concepto": "Dividendo provisorio", "monto_a": "$257,1", "monto_c": "$277,2", "monto_i": "$285,8"},
        {"fecha": "29-07-2025", "concepto": "Dividendo provisorio", "monto_a": "$118,9", "monto_c": "$139,1", "monto_i": "$147,7"},
        {"fecha": "22-10-2025", "concepto": "Dividendo provisorio", "monto_a": "$120,4", "monto_c": "$140,5", "monto_i": "$148,7"},
        {"fecha": "29-12-2025", "concepto": "Dividendo provisorio", "monto_a": "$191,2", "monto_c": "$211,5", "monto_i": "$219,8"}
      ],
      "otros_indicadores": {
        "tasa_arriendo_a": "9,3%", "tasa_arriendo_c": "9,3%", "tasa_arriendo_i": "7,8%",
        "cap_rate_a": "6,7%",      "cap_rate_c": "6,7%",      "cap_rate_i": "5,7%",
        "ingresos_u12m": "UF 392.322", "ingresos_mes": "UF 33.133",
        "noi_u12m": "UF 314.116",      "noi_mes": "UF 28.224"
      },
      "balance": {
        "fecha": "31 de diciembre de 2025",
        "efectivo": "...", "otros_activos_corrientes": "...",
        "propiedades_inversion": "...", "activo_impuestos_diferidos": "...", "total_activos": "...",
        "prestamos_bancarios": "...", "pasivos_impuestos_diferidos": "...",
        "otros_pasivos": "...", "patrimonio": "...", "total_pasivos_patrimonio": "..."
      },
      "gastos": {
        "fecha": "31 DE DICIEMBRE 2025",
        "comision": "...", "recurrentes": "...", "otros": "...", "total": "..."
      },
      "endeudamiento": {
        "leverage": "...", "ltv": "...", "tasa_promedio": "...", "duration": "...", "deuda_neta": "..."
      },
      "perfil_vencimiento": {"0_3": "...", "3_7": "...", "7_10": "...", "mas_10": "..."}
    }
    """
    try:
        from pptx import Presentation
    except ImportError:
        return "Error: python-pptx no instalado."

    work = _work_path("TRI")
    if not os.path.exists(work):
        return "No hay archivo de trabajo. Llama primero a preparar_fs('TRI', año, mes)."

    try:
        datos = json.loads(datos_json)
    except json.JSONDecodeError as e:
        return f"Error al parsear datos_json: {e}"

    prs = Presentation(work)
    slide = prs.slides[0]
    actualizados = []
    omitidos = []

    # Fecha del mes
    mes_str = f"{_MESES_ES[mes].upper()} {año}"
    updated_fecha = False
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            for m_es in _MESES_ES.values():
                if m_es.upper() in txt and any(str(y) in txt for y in range(2020, 2035)):
                    para = shape.text_frame.paragraphs[0]
                    if para.runs:
                        para.runs[0].text = mes_str
                        for run in para.runs[1:]:
                            run.text = ""
                    updated_fecha = True
                    break
        if updated_fecha:
            break
    if updated_fecha:
        actualizados.append(f"Fecha → {mes_str}")

    # Tabla 15: Valor cuota bursátil (3 meses × 3 series)
    pb = datos.get("precios_bursatiles", [])
    if pb:
        updates = {}
        for i, item in enumerate(pb[:3]):
            row = 2 + i
            updates[(row, 0)] = item.get("fecha", "")
            updates[(row, 1)] = item.get("serie_a", "")
            updates[(row, 2)] = item.get("serie_c", "")
            updates[(row, 3)] = item.get("serie_i", "")
        if _update_table(slide, "Tabla 15", {k: v for k, v in updates.items() if v}):
            actualizados.append("Tabla 15 (precios bursátiles)")
        else:
            omitidos.append("Tabla 15")

    # Tabla 3: Valor cuota libro (4 trimestres × 3 series)
    vl = datos.get("valor_libro", [])
    if vl:
        updates = {}
        for i, item in enumerate(vl[:4]):
            row = 2 + i
            updates[(row, 0)] = item.get("fecha", "")
            updates[(row, 1)] = item.get("serie_a", "")
            updates[(row, 2)] = item.get("serie_c", "")
            updates[(row, 3)] = item.get("serie_i", "")
        if _update_table(slide, "Tabla 3", {k: v for k, v in updates.items() if v}):
            actualizados.append("Tabla 3 (valor cuota libro)")
        else:
            omitidos.append("Tabla 3")

    # Tabla 11: Rentabilidad (multi-serie A/C/I × Bursátil/Libro)
    rent = datos.get("rentabilidad", {})
    if rent:
        # cols: 1=A_burs, 2=A_lib, 3=C_burs, 4=C_lib, 5=I_burs, 6=I_lib
        updates = {
            (3, 1): rent.get("inicio_bursatil_a", ""), (3, 2): rent.get("inicio_libro_a", ""),
            (3, 3): rent.get("inicio_bursatil_c", ""), (3, 4): rent.get("inicio_libro_c", ""),
            (3, 5): rent.get("inicio_bursatil_i", ""), (3, 6): rent.get("inicio_libro_i", ""),
            (4, 1): rent.get("ytd_bursatil_a", ""),    (4, 2): rent.get("ytd_libro_a", ""),
            (4, 3): rent.get("ytd_bursatil_c", ""),    (4, 4): rent.get("ytd_libro_c", ""),
            (4, 5): rent.get("ytd_bursatil_i", ""),    (4, 6): rent.get("ytd_libro_i", ""),
            (5, 1): rent.get("12m_bursatil_a", ""),    (5, 2): rent.get("12m_libro_a", ""),
            (5, 3): rent.get("12m_bursatil_c", ""),    (5, 4): rent.get("12m_libro_c", ""),
            (5, 5): rent.get("12m_bursatil_i", ""),    (5, 6): rent.get("12m_libro_i", ""),
            (6, 1): rent.get("dy_bursatil_a", ""),     (6, 2): rent.get("dy_libro_a", ""),
            (6, 3): rent.get("dy_bursatil_c", ""),     (6, 4): rent.get("dy_libro_c", ""),
            (6, 5): rent.get("dy_bursatil_i", ""),     (6, 6): rent.get("dy_libro_i", ""),
            (7, 1): rent.get("dy_amort_bursatil_a", ""),(7, 2): rent.get("dy_amort_libro_a", ""),
            (7, 3): rent.get("dy_amort_bursatil_c", ""),(7, 4): rent.get("dy_amort_libro_c", ""),
            (7, 5): rent.get("dy_amort_bursatil_i", ""),(7, 6): rent.get("dy_amort_libro_i", ""),
        }
        if _update_table(slide, "Tabla 11", {k: v for k, v in updates.items() if v}):
            actualizados.append("Tabla 11 (rentabilidad)")
        else:
            omitidos.append("Tabla 11")

    # Tabla 52: Repartos (4 dividendos × 3 series)
    divs = datos.get("dividendos", [])
    if divs:
        updates = {}
        for i, div in enumerate(divs[:4]):
            row = 2 + i
            updates[(row, 0)] = div.get("fecha", "")
            updates[(row, 1)] = div.get("concepto", "")
            updates[(row, 2)] = div.get("monto_a", "")
            updates[(row, 3)] = div.get("monto_c", "")
            updates[(row, 4)] = div.get("monto_i", "")
        if _update_table(slide, "Tabla 52", {k: v for k, v in updates.items() if v}):
            actualizados.append("Tabla 52 (repartos)")
        else:
            omitidos.append("Tabla 52")

    # Tabla 44: Otros indicadores (series en cols 1/2/3 para A/C/I; NOI/Ingresos solo col 1)
    noi = datos.get("otros_indicadores", {})
    if noi:
        updates = {
            (1, 1): noi.get("tasa_arriendo_a", ""),
            (1, 2): noi.get("tasa_arriendo_c", ""),
            (1, 3): noi.get("tasa_arriendo_i", ""),
            (2, 1): noi.get("cap_rate_a", ""),
            (2, 2): noi.get("cap_rate_c", ""),
            (2, 3): noi.get("cap_rate_i", ""),
            (3, 1): noi.get("ingresos_u12m", ""),
            (4, 1): noi.get("ingresos_mes", ""),
            (5, 1): noi.get("noi_u12m", ""),
            (6, 1): noi.get("noi_mes", ""),
        }
        if _update_table(slide, "Tabla 44", {k: v for k, v in updates.items() if v}):
            actualizados.append("Tabla 44 (otros indicadores)")
        else:
            omitidos.append("Tabla 44")

    # Tabla 5: Balance consolidado
    bal = datos.get("balance", {})
    if bal:
        fecha_bal = bal.get("fecha", "")
        updates = {}
        if fecha_bal:
            updates[(0, 0)] = f"BALANCE CONSOLIDADO AL {fecha_bal.upper()} (en miles de pesos)"
        updates.update({
            (1, 1): bal.get("efectivo", ""),
            (1, 4): bal.get("prestamos_bancarios", ""),
            (2, 1): bal.get("otros_activos_corrientes", ""),
            (2, 4): bal.get("pasivos_impuestos_diferidos", ""),
            (3, 1): bal.get("propiedades_inversion", ""),
            (3, 4): bal.get("otros_pasivos", ""),
            (4, 1): bal.get("activo_impuestos_diferidos", ""),
            (4, 4): bal.get("patrimonio", ""),
            (5, 1): bal.get("total_activos", ""),
            (5, 4): bal.get("total_pasivos_patrimonio", ""),
        })
        if _update_table(slide, "Tabla 5", {k: v for k, v in updates.items() if v}):
            actualizados.append("Tabla 5 (balance)")
        else:
            omitidos.append("Tabla 5")

    # Tabla 8: Gastos del fondo
    gas = datos.get("gastos", {})
    if gas:
        fecha_gas = gas.get("fecha", "")
        updates = {}
        if fecha_gas:
            updates[(0, 0)] = f"GASTOS DEL FONDO AL {fecha_gas.upper()} (en miles de pesos)"
        updates.update({
            (1, 1): gas.get("comision", ""),
            (2, 1): gas.get("recurrentes", ""),
            (3, 1): gas.get("otros", ""),
            (4, 1): gas.get("total", ""),
        })
        if _update_table(slide, "Tabla 8", {k: v for k, v in updates.items() if v}):
            actualizados.append("Tabla 8 (gastos)")
        else:
            omitidos.append("Tabla 8")

    # Tabla 38: Endeudamiento
    end = datos.get("endeudamiento", {})
    if end:
        updates = {
            (1, 1): end.get("leverage", ""),
            (2, 1): end.get("ltv", ""),
            (3, 1): end.get("tasa_promedio", ""),
            (4, 1): end.get("duration", ""),
            (5, 1): end.get("deuda_neta", ""),
        }
        if _update_table(slide, "Tabla 38", {k: v for k, v in updates.items() if v}):
            actualizados.append("Tabla 38 (endeudamiento)")
        else:
            omitidos.append("Tabla 38")

    # Tabla 2: Perfil vencimiento
    perf = datos.get("perfil_vencimiento", {})
    if perf:
        updates = {
            (1, 1): perf.get("0_3", ""),
            (2, 1): perf.get("3_7", ""),
            (3, 1): perf.get("7_10", ""),
            (4, 1): perf.get("mas_10", ""),
        }
        if _update_table(slide, "Tabla 2", {k: v for k, v in updates.items() if v}):
            actualizados.append("Tabla 2 (perfil vencimiento)")
        else:
            omitidos.append("Tabla 2")

    prs.save(work)

    resultado = f"FS TRI actualizado en WORK_DIR.\n"
    if actualizados:
        resultado += f"\nActualizados ({len(actualizados)}):\n" + "\n".join(f"  ✓ {x}" for x in actualizados)
    if omitidos:
        resultado += f"\nOmitidos:\n" + "\n".join(f"  ✗ {x}" for x in omitidos)
    resultado += f"\n\nLlama a guardar_fs('TRI', año, mes) para guardar en SharePoint."
    return resultado


def guardar_fs(fondo_key: str, año: int, mes: int) -> str:
    """
    Guarda el archivo de trabajo como nueva versión del Fact Sheet en SharePoint.
    Nomenclatura: YYMM Fact Sheet - <nombre fondo>.pptx
    """
    if fondo_key not in _FS_DIRS:
        return f"fondo_key inválido: {fondo_key}."

    work = _work_path(fondo_key)
    if not os.path.exists(work):
        return f"No hay archivo de trabajo para {fondo_key}. Llama primero a preparar_fs."

    yymm = f"{str(año)[-2:]}{mes:02d}"
    nombre = f"{yymm} {_FS_OUTPUT_BASE[fondo_key]}.pptx"

    dest_folder = _fs_folder(fondo_key, año, mes)
    os.makedirs(dest_folder, exist_ok=True)
    dest = os.path.join(dest_folder, nombre)

    shutil.copy2(work, dest)
    return (
        f"FS {fondo_key} guardado en SharePoint.\n"
        f"  Archivo: {nombre}\n"
        f"  Ruta:    {dest}"
    )
